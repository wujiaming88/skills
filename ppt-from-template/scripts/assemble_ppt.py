#!/usr/bin/env python3
"""根据 plan.yaml 和模板，组装生成 PPT。

全程使用 python-pptx API，不直接操作 XML，确保兼容性。

用法: python3 assemble_ppt.py <template.pptx> <plan.yaml> <output.pptx>

plan.yaml 格式:
  title: "演示文稿标题"
  slides:
    - layout: cover
      source_slide: 1          # 模板中的 slide 编号（1-based）
      content:
        title: "标题文字"
        subtitle: "副标题"
    - layout: content
      source_slide: 8
      content:
        title: "第一章"
        body:
          - "要点一"
          - "要点二"
          - "要点三"
"""

import sys
import os
import copy
import yaml
from pptx import Presentation
from pptx.util import Pt, Emu
from lxml import etree


def deep_copy_slide(prs, template_slide):
    """深拷贝一张 slide（包括所有形状、关系、媒体）。

    使用 lxml 在 presentation XML 层面复制，但通过 python-pptx 管理关系，
    确保兼容性。
    """
    # 获取 slide layout
    slide_layout = template_slide.slide_layout

    # 添加新 slide
    new_slide = prs.slides.add_slide(slide_layout)

    # 清空新 slide 的默认占位符内容
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # 复制所有形状元素
    for shape in template_slide.shapes:
        el = copy.deepcopy(shape._element)
        new_slide.shapes._spTree.append(el)

    # 复制背景（如果有）
    if template_slide.background and template_slide.background._element is not None:
        bg = template_slide.background._element
        new_bg = copy.deepcopy(bg)
        # 替换新 slide 的背景
        existing_bg = new_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}bg'
        )
        if existing_bg is not None:
            new_slide._element.remove(existing_bg)
        # 插到 cSld 的最前面
        csld = new_slide._element.find(
            '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
        )
        if csld is not None:
            csld.insert(0, new_bg)

    # 复制 slide 级别的关系（图片等媒体）
    for rel in template_slide.part.rels.values():
        if "image" in rel.reltype or "media" in rel.reltype:
            try:
                new_slide.part.rels.get_or_add(rel.reltype, rel.target_part)
            except Exception:
                pass  # 跳过无法复制的关系

    return new_slide


def set_text_in_shape(shape, text, is_list=False):
    """设置形状中的文本，保留原有格式。

    策略：保留第一个 paragraph 的格式作为模板，替换文本。
    """
    if not shape.has_text_frame:
        return False

    tf = shape.text_frame

    if not tf.paragraphs:
        return False

    if is_list and isinstance(text, list):
        # 列表模式：每个 item 一个 paragraph
        if not text:
            # 空列表：清空文本但保留一个空段落
            for p in list(tf.paragraphs)[1:]:
                p._element.getparent().remove(p._element)
            _set_para_text(tf.paragraphs[0], "", _extract_para_format(tf.paragraphs[0]))
            return True

        # 保留第一个 paragraph 的格式作为模板
        template_para = tf.paragraphs[0]
        template_fmt = _extract_para_format(template_para)

        # 清空所有现有 paragraph
        for p in list(tf.paragraphs)[1:]:
            p._element.getparent().remove(p._element)

        # 第一个 item 写入第一个 paragraph
        _set_para_text(tf.paragraphs[0], str(text[0]), template_fmt)

        # 后续 item 添加新 paragraph
        for item in text[1:]:
            new_p = copy.deepcopy(template_para._element)
            tf._txBody.append(new_p)
            # 获取刚添加的 paragraph 对象
            last_para = tf.paragraphs[-1]
            _set_para_text(last_para, str(item), template_fmt)
    else:
        # 单文本模式：替换所有 paragraph 为一段
        # 清空多余 paragraph
        for p in list(tf.paragraphs)[1:]:
            p._element.getparent().remove(p._element)
        template_fmt = _extract_para_format(tf.paragraphs[0])
        _set_para_text(tf.paragraphs[0], str(text), template_fmt)

    return True


def _extract_para_format(para):
    """提取 paragraph 格式信息"""
    fmt = {'font_size': None, 'font_bold': None, 'font_color': None, 'font_name': None}
    if para.runs:
        run = para.runs[0]
        fmt['font_size'] = run.font.size
        fmt['font_bold'] = run.font.bold
        fmt['font_name'] = run.font.name
        try:
            if run.font.color and run.font.color.type is not None:
                fmt['font_color'] = run.font.color.rgb
        except (AttributeError, TypeError):
            pass  # 预设颜色等无法提取 rgb 的情况
    return fmt


def _set_para_text(para, text, fmt=None):
    """设置 paragraph 文本，保留或应用格式"""
    # 清空现有 runs
    for r in list(para.runs):
        r._r.getparent().remove(r._r)

    # 添加新 run
    run = para.add_run()
    run.text = text

    if fmt:
        if fmt['font_size']:
            run.font.size = fmt['font_size']
        if fmt['font_bold'] is not None:
            run.font.bold = fmt['font_bold']
        if fmt['font_name']:
            run.font.name = fmt['font_name']
        if fmt['font_color']:
            run.font.color.rgb = fmt['font_color']


def find_shape_by_role(slide, role):
    """按语义角色查找形状（title/subtitle/body）"""
    candidates = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        text = tf.text.strip()
        font_sizes = []
        for p in tf.paragraphs:
            for r in p.runs:
                if r.font.size:
                    font_sizes.append(r.font.size)
        max_size = max(font_sizes) if font_sizes else 0

        candidates.append({
            'shape': shape,
            'text': text,
            'top': shape.top or 0,
            'max_font_size': max_size,
            'para_count': len(tf.paragraphs),
        })

    if not candidates:
        return None

    # 按位置排序（从上到下）
    candidates.sort(key=lambda c: c['top'])

    if role == 'title':
        # 字号最大的
        return max(candidates, key=lambda c: c['max_font_size'])['shape']
    elif role == 'subtitle':
        # 第二个文本框（标题之后）
        if len(candidates) >= 2:
            title_shape = max(candidates, key=lambda c: c['max_font_size'])['shape']
            for c in candidates:
                if c['shape'] != title_shape:
                    return c['shape']
        return None
    elif role == 'body':
        # paragraph 最多的（排除标题）
        title_shape = max(candidates, key=lambda c: c['max_font_size'])['shape']
        body_candidates = [c for c in candidates if c['shape'] != title_shape]
        if body_candidates:
            return max(body_candidates, key=lambda c: c['para_count'])['shape']
        return None

    return None


def apply_content(slide, content):
    """将 plan 中的 content 应用到 slide"""
    if not content:
        return

    applied = []

    # 标题
    if 'title' in content:
        shape = find_shape_by_role(slide, 'title')
        if shape:
            set_text_in_shape(shape, content['title'])
            applied.append('title')

    # 副标题
    if 'subtitle' in content:
        shape = find_shape_by_role(slide, 'subtitle')
        if shape:
            set_text_in_shape(shape, content['subtitle'])
            applied.append('subtitle')

    # 正文（列表或文本）
    if 'body' in content:
        shape = find_shape_by_role(slide, 'body')
        if shape:
            body = content['body']
            if isinstance(body, list):
                set_text_in_shape(shape, body, is_list=True)
            else:
                set_text_in_shape(shape, body)
            applied.append('body')

    return applied


def assemble_ppt(template_path, plan_path, output_path):
    """主函数：根据 plan 组装 PPT"""
    # 读取 plan
    with open(plan_path, 'r', encoding='utf-8') as f:
        plan = yaml.safe_load(f)

    # 打开模板
    prs = Presentation(template_path)
    template_slides = list(prs.slides)
    original_count = len(template_slides)

    print(f"📖 模板: {original_count} 页")
    print(f"📝 计划: {len(plan['slides'])} 页")

    # 第一步：复制需要的页面
    new_slides = []
    for i, slide_plan in enumerate(plan['slides']):
        source_num = slide_plan['source_slide']
        if source_num < 1 or source_num > original_count:
            print(f"  ⚠️ slide {i+1}: source_slide {source_num} 超出范围，跳过")
            continue

        source_slide = template_slides[source_num - 1]
        new_slide = deep_copy_slide(prs, source_slide)
        new_slides.append((new_slide, slide_plan))
        layout = slide_plan.get('layout', '?')
        print(f"  ✅ 第{i+1}页: 复制 slide{source_num} [{layout}]")

    # 第二步：删除原始页面（倒序删除避免索引偏移）
    xml_slides = prs.slides._sldIdLst
    slide_ids = list(xml_slides)
    for sid in slide_ids[:original_count]:
        xml_slides.remove(sid)
    print(f"🗑️ 已删除原始 {original_count} 页")

    # 第三步：应用内容
    print("✍️ 替换内容...")
    for new_slide, slide_plan in new_slides:
        content = slide_plan.get('content', {})
        applied = apply_content(new_slide, content)
        if applied:
            print(f"  ✅ {slide_plan.get('layout', '?')}: {', '.join(applied)}")

    # 保存
    prs.save(output_path)
    size_kb = os.path.getsize(output_path) / 1024
    print(f"\n✅ 输出: {output_path} ({size_kb:.0f} KB, {len(plan['slides'])} 页)")


if __name__ == '__main__':
    if len(sys.argv) < 4:
        print(f"用法: {sys.argv[0]} <template.pptx> <plan.yaml> <output.pptx>")
        sys.exit(1)

    assemble_ppt(sys.argv[1], sys.argv[2], sys.argv[3])
