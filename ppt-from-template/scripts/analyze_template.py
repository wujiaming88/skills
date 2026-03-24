#!/usr/bin/env python3
"""分析 PPT 模板，提取版式信息，输出 layouts.yaml。

用法: python3 analyze_template.py <template.pptx> [output_dir]

输出:
  - layouts.yaml: 版式定义（每种版式的 slide 编号、占位符、图片位置等）
  - thumbnails/slideN.txt: 每页的文本摘要（辅助人工确认版式分类）
"""

import sys
import os
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu


def emu_to_cm(emu):
    """EMU 转厘米，保留 1 位小数"""
    if emu is None:
        return None
    return round(emu / 914400 * 2.54, 1)


def classify_layout(shapes_info):
    """根据形状组合推断版式类型"""
    has_title = any(s['role'] == 'title' for s in shapes_info)
    has_subtitle = any(s['role'] == 'subtitle' for s in shapes_info)
    has_body = any(s['role'] == 'body' for s in shapes_info)
    has_image = any(s['type'] == 'picture' for s in shapes_info)
    has_table = any(s['type'] == 'table' for s in shapes_info)
    has_chart = any(s['type'] == 'chart' for s in shapes_info)
    text_count = sum(1 for s in shapes_info if s['type'] == 'text' and s['role'] not in ('title', 'subtitle'))

    if has_title and not has_body and not has_image and text_count <= 2:
        if has_subtitle:
            return 'cover'        # 封面：标题 + 副标题
        return 'section'          # 章节页：仅标题
    if has_table or has_chart:
        return 'data'             # 数据页
    if has_image and has_body:
        return 'image-text'       # 图文页
    if has_image and not has_body:
        return 'image'            # 纯图片页
    if has_body or text_count >= 2:
        return 'content'          # 内容页（标题 + 要点）
    return 'other'                # 其他


def guess_role(shape, idx, all_shapes):
    """推断形状的语义角色"""
    if shape.has_text_frame:
        tf = shape.text_frame
        text = tf.text.strip()
        # 位置最高 + 字号最大 → 标题
        font_sizes = []
        for para in tf.paragraphs:
            for run in para.runs:
                if run.font.size:
                    font_sizes.append(run.font.size)
        max_size = max(font_sizes) if font_sizes else 0

        # 标题启发：位于上半部分且字号 ≥ 24pt
        if shape.top is not None and shape.top < Emu(2500000) and max_size >= Pt(20):
            return 'title'
        # 副标题：紧接标题下方，字号中等
        if shape.top is not None and Emu(2500000) <= shape.top < Emu(4000000) and Pt(12) <= max_size < Pt(20):
            return 'subtitle'
        if len(tf.paragraphs) >= 3:
            return 'body'         # 多行 → 正文
        return 'text'
    return 'other'


def analyze_shape(shape, idx, all_shapes):
    """提取单个形状的信息"""
    info = {
        'index': idx,
        'name': shape.name,
        'type': 'other',
        'role': 'other',
        'position': {
            'left_cm': emu_to_cm(shape.left),
            'top_cm': emu_to_cm(shape.top),
            'width_cm': emu_to_cm(shape.width),
            'height_cm': emu_to_cm(shape.height),
        },
        'text_preview': '',
    }

    if shape.has_text_frame:
        info['type'] = 'text'
        info['role'] = guess_role(shape, idx, all_shapes)
        text = shape.text_frame.text.strip()
        info['text_preview'] = text[:80] + ('...' if len(text) > 80 else '')
        # 提取段落数和字号
        paras = shape.text_frame.paragraphs
        info['paragraph_count'] = len(paras)
        sizes = set()
        for p in paras:
            for r in p.runs:
                if r.font.size:
                    sizes.add(round(r.font.size / Pt(1), 1))
        info['font_sizes_pt'] = sorted(sizes) if sizes else []
    elif shape.shape_type is not None:
        type_name = str(shape.shape_type)
        if 'PICTURE' in type_name or 'IMAGE' in type_name:
            info['type'] = 'picture'
        elif 'TABLE' in type_name:
            info['type'] = 'table'
        elif 'CHART' in type_name:
            info['type'] = 'chart'
        elif 'GROUP' in type_name:
            info['type'] = 'group'
        elif 'PLACEHOLDER' in type_name:
            info['type'] = 'placeholder'
            if shape.has_text_frame:
                info['role'] = guess_role(shape, idx, all_shapes)
                info['text_preview'] = shape.text_frame.text.strip()[:80]

    return info


def analyze_template(pptx_path, output_dir):
    """分析模板并输出 layouts.yaml"""
    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)
    os.makedirs(os.path.join(output_dir, 'thumbnails'), exist_ok=True)

    slides_info = []
    layout_groups = {}  # layout_type -> [slide_numbers]

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        shapes_info = []
        for shape_idx, shape in enumerate(slide.shapes):
            si = analyze_shape(shape, shape_idx, slide.shapes)
            shapes_info.append(si)

        layout_type = classify_layout(shapes_info)

        slide_data = {
            'slide_number': slide_num,
            'layout_type': layout_type,
            'shape_count': len(shapes_info),
            'shapes': shapes_info,
        }
        slides_info.append(slide_data)

        # 按版式类型分组
        if layout_type not in layout_groups:
            layout_groups[layout_type] = []
        layout_groups[layout_type].append(slide_num)

        # 写文本摘要
        text_lines = [f"Slide {slide_num} [{layout_type}]"]
        for si in shapes_info:
            if si['text_preview']:
                text_lines.append(f"  [{si['role']}] {si['text_preview']}")
        thumb_path = os.path.join(output_dir, 'thumbnails', f'slide{slide_num}.txt')
        with open(thumb_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(text_lines))

    # 为每种版式选取代表 slide
    layouts = {}
    for lt, slide_nums in layout_groups.items():
        # 取第一个作为代表
        representative = slide_nums[0]
        rep_data = slides_info[representative - 1]

        # 提取可编辑占位符
        placeholders = []
        for s in rep_data['shapes']:
            if s['type'] == 'text' and s['role'] in ('title', 'subtitle', 'body', 'text'):
                placeholders.append({
                    'shape_index': s['index'],
                    'role': s['role'],
                    'text_preview': s['text_preview'],
                    'position': s['position'],
                })

        layouts[lt] = {
            'description': _layout_description(lt),
            'representative_slide': representative,
            'all_slides': slide_nums,
            'placeholders': placeholders,
        }

    # 输出 layouts.yaml
    output = {
        'template': os.path.basename(pptx_path),
        'total_slides': len(prs.slides),
        'slide_width_cm': emu_to_cm(prs.slide_width),
        'slide_height_cm': emu_to_cm(prs.slide_height),
        'layouts': layouts,
        'slides_detail': slides_info,
    }

    yaml_path = os.path.join(output_dir, 'layouts.yaml')
    with open(yaml_path, 'w', encoding='utf-8') as f:
        yaml.dump(output, f, allow_unicode=True, default_flow_style=False, sort_keys=False)

    print(f"✅ 分析完成: {len(prs.slides)} 页, {len(layouts)} 种版式")
    print(f"   版式: {', '.join(f'{k}({len(v['all_slides'])}页)' for k, v in layouts.items())}")
    print(f"   输出: {yaml_path}")
    return output


def _layout_description(lt):
    descs = {
        'cover': '封面 - 大标题 + 副标题',
        'section': '章节页 - 章节标题',
        'content': '内容页 - 标题 + 多行要点',
        'image-text': '图文页 - 图片 + 文字说明',
        'image': '纯图片页',
        'data': '数据页 - 表格或图表',
        'other': '其他',
    }
    return descs.get(lt, lt)


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(f"用法: {sys.argv[0]} <template.pptx> [output_dir]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else os.path.dirname(pptx_path) or '.'
    analyze_template(pptx_path, output_dir)
